"""Step 9: PPTX extraction — per-slide sectioning with slide-number ranges, speaker notes,
tables, core properties, fallbacks.

Synthetic .pptx files built with python-pptx so the tests are deterministic and committable.
"""

from datetime import datetime
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from locus.extract.base import MAX_SECTION_CHARS
from locus.extract.pptx import extract_pptx

BODY = " ".join(f"Sentence {i} contains several descriptive words about the topic." for i in range(12))

# Default-template layout indices: 5 = Title Only, 6 = Blank.
_TITLE_ONLY, _BLANK = 5, 6


def _add_slide(prs, *, title=None, body=None, notes=None, layout=_TITLE_ONLY):
    slide = prs.slides.add_slide(prs.slide_layouts[layout])
    if title is not None:
        slide.shapes.title.text = title
    if body is not None:
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        box.text_frame.text = body
    if notes is not None:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def _save(prs, path: Path) -> Path:
    prs.save(str(path))
    return path


def test_per_slide_sections_carry_slide_numbers(tmp_path: Path):
    prs = Presentation()
    prs.core_properties.title = "Quarterly Architecture Review"
    prs.core_properties.created = datetime(2024, 3, 9, 12, 0, 0)
    for name in ("Background", "Implementation", "Results"):
        _add_slide(prs, title=name, body=BODY)
    doc = extract_pptx(_save(prs, tmp_path / "review.pptx"))
    assert doc.title == "Quarterly Architecture Review"
    assert doc.source_date == "2024-03-09"
    assert doc.section_strategy == "slides"
    assert doc.page_count == 3
    assert [s.title for s in doc.sections] == ["Background", "Implementation", "Results"]
    assert [(s.page_start, s.page_end) for s in doc.sections] == [(1, 1), (2, 2), (3, 3)]
    assert doc.sections[0].text.startswith("Background")  # title kept in verbatim text


def test_speaker_notes_appended(tmp_path: Path):
    prs = Presentation()
    _add_slide(prs, title="Roadmap", body=BODY, notes="Mention the Q3 dependency risk.")
    doc = extract_pptx(_save(prs, tmp_path / "notes.pptx"))
    text = doc.sections[0].text
    assert "Notes:" in text
    assert "Mention the Q3 dependency risk." in text


def test_table_text_is_retrievable(tmp_path: Path):
    prs = Presentation()
    slide = _add_slide(prs, title="Results", body=BODY)
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Sharpe ratio"
    table.cell(1, 1).text = "1.42"
    doc = extract_pptx(_save(prs, tmp_path / "table.pptx"))
    joined = "\n".join(s.text for s in doc.sections)
    assert "Sharpe ratio | 1.42" in joined


def test_small_slides_merge_keeping_slide_range(tmp_path: Path):
    prs = Presentation()
    for i in range(4):
        _add_slide(prs, title=f"Point {i + 1}", body="One short bullet.")
    doc = extract_pptx(_save(prs, tmp_path / "tiny.pptx"))
    assert len(doc.sections) == 1
    assert doc.sections[0].page_start == 1
    assert doc.sections[0].page_end == 4
    assert doc.sections[0].title == "Point 1"  # first non-empty heading survives the merge
    assert "Point 4" in doc.sections[0].text  # absorbed headings stay visible


def test_oversized_slide_windows_inherit_range(tmp_path: Path):
    prs = Presentation()
    long_notes = "\n\n".join(BODY for _ in range(20))  # well over MAX_SECTION_CHARS
    _add_slide(prs, title="Dense Slide", body=BODY, notes=long_notes)
    doc = extract_pptx(_save(prs, tmp_path / "dense.pptx"))
    assert len(doc.sections) > 1
    assert all(len(s.text) <= MAX_SECTION_CHARS for s in doc.sections)
    assert all((s.page_start, s.page_end) == (1, 1) for s in doc.sections)
    assert doc.sections[0].title == "Dense Slide (part 1)"


def test_chrome_placeholders_filtered(tmp_path: Path):
    # PowerPoint injects slide-number/footer/date placeholders into slide XML when
    # "Insert > Header & Footer" is on; add_slide does NOT clone them from the layout,
    # so replicate PowerPoint by deep-copying the layout's chrome placeholders in.
    import copy

    from pptx.enum.shapes import PP_PLACEHOLDER

    prs = Presentation()
    layout = prs.slide_layouts[_TITLE_ONLY]
    slide = _add_slide(prs, title="Findings", body=BODY, layout=_TITLE_ONLY)
    for ph in layout.placeholders:
        if ph.placeholder_format.type in (
            PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.DATE,
        ):
            slide.shapes._spTree.append(copy.deepcopy(ph._element))
    chrome = {
        PP_PLACEHOLDER.SLIDE_NUMBER: "7",
        PP_PLACEHOLDER.FOOTER: "Company Confidential",
        PP_PLACEHOLDER.DATE: "2026-06-05",
    }
    for ph in slide.placeholders:
        kind = ph.placeholder_format.type
        if kind in chrome:
            ph.text_frame.text = chrome[kind]

    doc = extract_pptx(_save(prs, tmp_path / "chrome.pptx"))
    text = doc.sections[0].text
    assert text.startswith("Findings")  # title survives, not the slide number
    assert "Company Confidential" not in text
    assert "2026-06-05" not in text
    assert "\n7\n" not in f"\n{text}\n"  # the bare slide number is gone


def test_literal_slide_number_textbox_filtered(tmp_path: Path):
    # Google-Slides-style exports render page numbers as plain TEXT_BOX shapes whose
    # whole text is the slide number (no placeholder). Slide 1's own number is dropped;
    # a bare number that is NOT the slide's own number is content and stays.
    prs = Presentation()
    slide = _add_slide(prs, title="Scores", body=BODY)
    num = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(0.5), Inches(0.3))
    num.text_frame.text = "1"  # == its slide number -> chrome
    other = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(0.5), Inches(0.3))
    other.text_frame.text = "42"  # bare number, but not the slide number -> content
    doc = extract_pptx(_save(prs, tmp_path / "gslides.pptx"))
    text = doc.sections[0].text
    assert "\n1\n" not in f"\n{text}\n"
    assert "42" in text


def test_title_only_slide_emits_section(tmp_path: Path):
    prs = Presentation()
    _add_slide(prs, title="Closing Thoughts")
    doc = extract_pptx(_save(prs, tmp_path / "titleonly.pptx"))
    assert len(doc.sections) == 1
    assert doc.sections[0].text == "Closing Thoughts"


def test_title_falls_back_to_first_slide_title_then_stem(tmp_path: Path):
    prs = Presentation()  # no core-properties title
    _add_slide(prs, title="Actual Deck Heading", body=BODY)
    doc = extract_pptx(_save(prs, tmp_path / "untitled.pptx"))
    assert doc.title == "Actual Deck Heading"

    prs2 = Presentation()
    _add_slide(prs2, body=BODY, layout=_BLANK)
    doc2 = extract_pptx(_save(prs2, tmp_path / "plain-slides.pptx"))
    assert doc2.title == "plain-slides"
    assert doc2.section_strategy == "single"


def test_empty_deck_raises(tmp_path: Path):
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[_BLANK])
    with pytest.raises(ValueError, match="no extractable text"):
        extract_pptx(_save(prs, tmp_path / "empty.pptx"))


def test_non_pptx_raises_value_error(tmp_path: Path):
    bogus = tmp_path / "bogus.pptx"
    bogus.write_bytes(b"this is not a zip archive")
    with pytest.raises(ValueError, match="not a readable pptx"):
        extract_pptx(bogus)
