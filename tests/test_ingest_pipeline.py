"""Stage 5: ingest orchestration — write, idempotency, quarantine, routing.

The LLM passes and embeddings are monkeypatched with deterministic fakes, so these tests are
fast and need no Ollama. They exercise the orchestration and the transactional DB write.
"""

from pathlib import Path

import pymupdf
import pytest

from locus.db.connection import get_connection
from locus.db.migrate import migrate
from locus.ingest import embed, entities, gaps, propositions, summarize, synthesis
from locus.ingest.entities import Entity
from locus.ingest.llm import IngestExtractionError
from locus.ingest.synthesis import DocSynthesis
from locus import ingest_pipeline
from locus.ingest_pipeline import delete_document, ingest_file


def _make_pdf(path: Path) -> Path:
    doc = pymupdf.open()
    page = doc.new_page()
    y = 72.0
    lines = [("Section A", 20.0)]
    lines += [(f"Section A body sentence {i} with several descriptive words.", 11.0) for i in range(10)]
    lines += [("Section B", 20.0)]
    lines += [(f"Section B body sentence {i} describing methods and results.", 11.0) for i in range(10)]
    for text, size in lines:
        page.insert_text((72, y), text, fontsize=size)
        y += size + 6
    doc.save(str(path))
    doc.close()
    return path


def _count(conn, table: str) -> int:
    return conn.execute(f"SELECT COUNT(*) AS c FROM {table}").fetchone()["c"]


@pytest.fixture()
def conn(tmp_path: Path):
    db = tmp_path / "ingest.db"
    migrate(db)
    c = get_connection(db)
    yield c
    c.close()


@pytest.fixture()
def fake_passes(monkeypatch):
    """Replace the LLM passes + embeddings with deterministic stand-ins."""
    monkeypatch.setattr(
        summarize, "summarize_section",
        lambda title, text, **k: summarize.SectionSummary(
            summary=f"summary::{title}", title="Semantic Title"
        ),
    )
    monkeypatch.setattr(propositions, "extract_propositions", lambda title, text, **k: ["claim A", "claim B"])
    monkeypatch.setattr(entities, "extract_entities", lambda title, text, **k: [Entity(name="Kalman filter", type="method")])
    monkeypatch.setattr(
        synthesis, "synthesize_document",
        lambda title, summaries, **k: DocSynthesis(thesis="T", method="M", result="R", limitations="L"),
    )
    monkeypatch.setattr(gaps, "flag_gaps", lambda title, context, **k: ["gap one"])
    monkeypatch.setattr(embed, "embed_text", lambda text: [0.1] * 768)
    monkeypatch.setattr(embed, "embed_texts", lambda texts: [[0.1] * 768 for _ in texts])
    # Do NOT write synthetic test PDFs into the real vault/raw store (test isolation).
    monkeypatch.setattr(
        ingest_pipeline, "_copy_to_raw", lambda path, content_hash: f"{content_hash}{Path(path).suffix}"
    )


def test_failed_optional_pass_degrades_not_quarantines(tmp_path, conn, fake_passes, monkeypatch):
    """A section whose propositions pass persistently fails must not quarantine the doc:
    raw chunks + summary still carry the value (§15.0); the failure lands in gap_flags."""
    def boom(title, text, **k):
        raise IngestExtractionError("Propositions: no schema-valid output")
    monkeypatch.setattr(propositions, "extract_propositions", boom)

    result = ingest_file(_make_pdf(tmp_path / "doc.pdf"), conn)
    assert result.status == "ingested"  # not quarantined
    assert result.propositions == 0
    assert result.sections >= 2 and result.chunks >= 1  # everything else intact
    import json
    flags = json.loads(conn.execute("SELECT gap_flags FROM documents").fetchone()["gap_flags"])
    assert any("propositions pass failed" in f for f in flags)


def test_ingest_populates_all_levels(tmp_path, conn, fake_passes):
    result = ingest_file(_make_pdf(tmp_path / "doc.pdf"), conn)

    assert result.status == "ingested"
    assert result.doc_id is not None
    assert result.sections >= 2

    assert _count(conn, "documents") == 1
    doc = conn.execute("SELECT thesis, gap_flags FROM documents").fetchone()
    assert doc["thesis"] == "T"
    import json
    assert json.loads(doc["gap_flags"]) == ["gap one"]

    # Every level + its vectors written, counts consistent.
    assert _count(conn, "sections") == result.sections
    assert _count(conn, "section_vectors") == result.sections
    assert _count(conn, "chunks") == result.chunks > 0
    assert _count(conn, "chunk_vectors") == result.chunks
    assert _count(conn, "propositions") == result.propositions == result.sections * 2
    assert _count(conn, "proposition_vectors") == result.propositions
    assert _count(conn, "entities") == result.entities >= 1


def test_source_date_from_pdf_metadata_and_category_from_folder(tmp_path, conn, fake_passes):
    # A PDF carrying a creation date, dropped under a 'papers' folder.
    papers = tmp_path / "papers"
    papers.mkdir()
    pdf = _make_pdf(papers / "doc.pdf")
    doc = pymupdf.open(str(pdf))
    doc.set_metadata({"creationDate": "D:20230115120000+00'00'"})
    doc.saveIncr()
    doc.close()

    ingest_file(pdf, conn)
    row = conn.execute("SELECT source_date, category FROM documents").fetchone()
    assert row["source_date"] == "2023-01-15"  # parsed from PDF metadata
    assert row["category"] == "paper"          # papers/ -> paper (drop-folder convention)


def test_source_date_falls_back_to_mtime(tmp_path, conn, fake_passes):
    # No embedded date and no category folder: mtime fallback, 'uncategorized'.
    from datetime import datetime, timezone

    pdf = _make_pdf(tmp_path / "doc.pdf")
    ingest_file(pdf, conn)
    row = conn.execute("SELECT source_date, category FROM documents").fetchone()
    expected = datetime.fromtimestamp(pdf.stat().st_mtime, tz=timezone.utc).date().isoformat()
    assert row["source_date"] == expected
    assert row["category"] == "uncategorized"


def test_reingest_is_a_noop(tmp_path, conn, fake_passes):
    pdf = _make_pdf(tmp_path / "doc.pdf")
    first = ingest_file(pdf, conn)
    docs_after_first = _count(conn, "documents")
    chunks_after_first = _count(conn, "chunks")

    second = ingest_file(pdf, conn)
    assert second.status == "skipped"
    assert second.doc_id == first.doc_id
    assert _count(conn, "documents") == docs_after_first == 1
    assert _count(conn, "chunks") == chunks_after_first


def test_failure_is_quarantined_with_no_partial_write(tmp_path, conn, monkeypatch, fake_passes):
    # A pass blows up mid-prepare; nothing must be written.
    def boom(*a, **k):
        raise IngestExtractionError("model produced garbage")

    monkeypatch.setattr(summarize, "summarize_section", boom)

    result = ingest_file(_make_pdf(tmp_path / "doc.pdf"), conn)
    assert result.status == "quarantined"
    assert "garbage" in (result.error or "")
    assert _count(conn, "documents") == 0
    assert _count(conn, "chunks") == 0


def test_pseudo_titles_get_semantic_replacement():
    from locus.ingest_pipeline import _needs_semantic_title

    # Pagination fallbacks and missing titles are replaced (2026-06-05 evaluation).
    assert _needs_semantic_title(None)
    assert _needs_semantic_title("Section (pp 6–9)")
    assert _needs_semantic_title("Section (pp 10-15)")  # ASCII hyphen variant
    # Real extractor headings — including split ones carrying a page suffix — are kept.
    assert not _needs_semantic_title("Boundary conditions")
    assert not _needs_semantic_title("Stability (pp 6–9)")
    assert not _needs_semantic_title("3 Methods")


def test_semantic_title_flows_into_written_section(tmp_path, conn, fake_passes, monkeypatch):
    # Force the extractor to produce an untitled (single/paginated) doc, then check the
    # summary pass's title lands in the sections table.
    import locus.extract.pdf as pdf_extract

    real_extract = pdf_extract.extract_pdf

    def no_headings(path, **kw):
        doc = real_extract(path, **kw)
        for s in doc.sections:
            s.title = None
        return doc

    monkeypatch.setattr(ingest_pipeline.pdf_extract, "extract_pdf", no_headings)
    result = ingest_file(_make_pdf(tmp_path / "doc.pdf"), conn)
    assert result.status == "ingested"
    titles = [r["title"] for r in conn.execute("SELECT title FROM sections")]
    assert all(t == "Semantic Title" for t in titles)  # from the fake summarize pass


def test_synthesis_title_overrides_extractor_title(tmp_path, conn, fake_passes, monkeypatch):
    # The synthesis pass arbitrates the doc title: a corrected title replaces the extractor's
    # heuristic pick ("ENGINEERING SCIENCE" banner-grab class of failure).
    monkeypatch.setattr(
        synthesis, "synthesize_document",
        lambda title, summaries, **k: DocSynthesis(
            thesis="T", method="M", result="R", limitations="L",
            title="2nd Year Syllabus 2025-2026",
        ),
    )
    result = ingest_file(_make_pdf(tmp_path / "doc.pdf"), conn)
    assert result.status == "ingested"
    row = conn.execute("SELECT title FROM documents").fetchone()
    assert row["title"] == "2nd Year Syllabus 2025-2026"


def test_unsupported_type_is_reported(tmp_path, conn):
    # .txt became a supported type in step 8; .bin remains genuinely unsupported.
    blob = tmp_path / "firmware.bin"
    blob.write_text("not an ingestible format")
    result = ingest_file(blob, conn)
    assert result.status == "unsupported"


def test_source_type_routing_by_suffix():
    from locus.ingest_pipeline import _source_type

    cases = {
        "a.pdf": "pdf", "b.docx": "docx", "c.md": "markdown", "d.markdown": "markdown",
        "e.txt": "text", "f.ipynb": "notebook", "g.pptx": "slides",
        "H.MD": "markdown", "i.bin": None,
    }
    for name, expected in cases.items():
        assert _source_type(Path(name)) == expected, name


def test_markdown_ingests_end_to_end(tmp_path, conn, fake_passes):
    body = " ".join(f"Sentence {i} contains several descriptive words about the topic." for i in range(12))
    md = tmp_path / "note.md"
    md.write_text(
        f"---\ntitle: Control Systems Field Notes\ndate: 2023-07-21\n---\n\n"
        f"# Overview\n\n{body}\n\n# Details\n\n{body}\n",
        encoding="utf-8",
    )
    result = ingest_file(md, conn)
    assert result.status == "ingested"
    assert result.sections == 2
    row = conn.execute(
        "SELECT source_type, title, source_date FROM documents"
    ).fetchone()
    assert row["source_type"] == "markdown"
    assert row["title"] == "Control Systems Field Notes"  # trusted frontmatter title kept
    assert row["source_date"] == "2023-07-21"  # frontmatter date beats mtime fallback


def test_slides_ingest_end_to_end(tmp_path, conn, fake_passes):
    from pptx import Presentation
    from pptx.util import Inches

    body = " ".join(f"Sentence {i} contains several descriptive words about the topic." for i in range(12))
    prs = Presentation()
    for name in ("Background", "Results"):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        slide.shapes.title.text = name
        box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        box.text_frame.text = body
    deck = tmp_path / "deck.pptx"
    prs.save(str(deck))

    result = ingest_file(deck, conn)
    assert result.status == "ingested"
    assert result.sections == 2
    # Default pass profile (unlike code): propositions ran for every section.
    assert result.propositions == result.sections * 2

    import json
    row = conn.execute("SELECT source_type, section_map FROM documents").fetchone()
    assert row["source_type"] == "slides"
    # section_map carries REAL slide-number ranges, not the unpaginated 1/1.
    section_map = json.loads(row["section_map"])
    assert [(s["page_start"], s["page_end"]) for s in section_map] == [(1, 1), (2, 2)]


def test_missing_file_is_quarantined(tmp_path, conn):
    result = ingest_file(tmp_path / "nope.pdf", conn)
    assert result.status == "quarantined"


def test_delete_document_clears_vectors_too(tmp_path, conn, fake_passes):
    ingest_file(_make_pdf(tmp_path / "doc.pdf"), conn)
    assert _count(conn, "chunk_vectors") > 0

    doc_id = conn.execute("SELECT id FROM documents").fetchone()["id"]
    delete_document(conn, doc_id)

    # Both the FK-cascaded tables and the vec0 tables (no FK) must be empty.
    for t in ("documents", "sections", "chunks", "propositions", "entities",
              "section_vectors", "chunk_vectors", "proposition_vectors"):
        assert _count(conn, t) == 0, f"{t} not cleared"


def test_reingest_rebuilds_without_orphan_vectors(tmp_path, conn, fake_passes):
    pdf = _make_pdf(tmp_path / "doc.pdf")
    ingest_file(pdf, conn)
    first_chunks = _count(conn, "chunks")

    again = ingest_file(pdf, conn, reingest=True)
    assert again.status == "ingested"  # rebuilt, not skipped
    assert _count(conn, "documents") == 1
    # vectors stay exactly in step with their parents (no orphans accumulated).
    assert _count(conn, "chunk_vectors") == _count(conn, "chunks") == first_chunks
    assert _count(conn, "proposition_vectors") == _count(conn, "propositions")
    assert _count(conn, "section_vectors") == _count(conn, "sections")


# --- figures through the pipeline (step 11) -------------------------------------------------


def _make_figure_pdf(path: Path, *, caption: str = "Figure 3: Closed-loop control system.") -> Path:
    """A PDF with body text, one vector diagram, and a caption."""
    doc = pymupdf.open()
    page = doc.new_page()
    y = 72.0
    lines = [("Section A", 20.0)]
    lines += [(f"Section A body sentence {i} with several descriptive words.", 11.0) for i in range(12)]
    for text, size in lines:
        page.insert_text((72, y), text, fontsize=size)
        y += size + 6
    rect = pymupdf.Rect(100, 420, 450, 660)
    for i in range(12):
        yy = 420 + i * 20
        page.draw_line((100, yy), (450, yy - 10))
    page.draw_rect(rect)
    page.insert_text((100, 690), caption, fontsize=10)
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture()
def fig_env(tmp_path, fake_passes, monkeypatch):
    """Isolate the raw store and fake the VLM; returns the call log."""
    from locus.config import load as cfg_load
    from locus.ingest import figures as figures_mod

    monkeypatch.setattr(cfg_load().paths, "raw_store", tmp_path / "raw")
    calls: list[str] = []

    def fake_describe(image_bytes, caption=None, **kw):
        calls.append(caption or "")
        return "A closed-loop block diagram with controller and plant connected by feedback."

    monkeypatch.setattr(figures_mod, "describe_figure", fake_describe)
    return calls


def test_figures_written_with_vector_and_raw_png(tmp_path, conn, fig_env):
    from locus.config import load as cfg_load

    result = ingest_file(_make_figure_pdf(tmp_path / "doc.pdf"), conn)
    assert result.status == "ingested"
    assert result.figures == 1
    row = conn.execute("SELECT * FROM figures").fetchone()
    assert row["kind"] == "vector"
    assert row["caption"].startswith("Figure 3")
    assert row["description"].startswith("A closed-loop block diagram")
    assert row["section_id"] is not None
    assert _count(conn, "figure_vectors") == 1
    png = cfg_load().paths.raw_store / row["raw_path"]
    assert png.exists() and png.read_bytes()[:8] == b"\x89PNG\r\n\x1a\n"
    import json

    flags = json.loads(conn.execute("SELECT gap_flags FROM documents").fetchone()["gap_flags"])
    assert not any("figure description failed" in f for f in flags)


def test_reingest_hits_description_cache(tmp_path, conn, fig_env):
    pdf = _make_figure_pdf(tmp_path / "doc.pdf")
    ingest_file(pdf, conn)
    assert len(fig_env) == 1  # one VLM call on first ingest
    result = ingest_file(pdf, conn, reingest=True)
    assert result.status == "ingested" and result.figures == 1
    assert len(fig_env) == 1  # cache hit: NO second VLM call
    assert _count(conn, "figures") == 1 and _count(conn, "figure_vectors") == 1


def test_reingest_removes_orphaned_figure_pngs(tmp_path, conn, fig_env, monkeypatch):
    from locus.config import load as cfg_load
    from locus.extract import figures_detect

    pdf = _make_figure_pdf(tmp_path / "doc.pdf")
    ingest_file(pdf, conn)
    raw = cfg_load().paths.raw_store
    first = conn.execute("SELECT raw_path FROM figures").fetchone()["raw_path"]
    assert (raw / first).exists()

    # Detection rules changed (e.g. stricter filters): same content now yields no figures.
    monkeypatch.setattr(figures_detect, "detect_figures", lambda *a, **k: [])
    result = ingest_file(pdf, conn, reingest=True)
    assert result.status == "ingested" and result.figures == 0
    assert _count(conn, "figures") == 0 and _count(conn, "figure_vectors") == 0
    assert not (raw / first).exists()  # orphan removed after commit


def test_delete_document_removes_figure_pngs(tmp_path, conn, fig_env):
    from locus.config import load as cfg_load

    result = ingest_file(_make_figure_pdf(tmp_path / "doc.pdf"), conn)
    raw_path = conn.execute("SELECT raw_path FROM figures").fetchone()["raw_path"]
    png = cfg_load().paths.raw_store / raw_path
    assert png.exists()
    delete_document(conn, result.doc_id)
    assert _count(conn, "figures") == 0 and _count(conn, "figure_vectors") == 0
    assert not png.exists()


def test_failed_description_stores_caption_only(tmp_path, conn, fig_env, monkeypatch):
    from locus.ingest import figures as figures_mod

    monkeypatch.setattr(figures_mod, "describe_figure", lambda *a, **k: None)
    result = ingest_file(_make_figure_pdf(tmp_path / "doc.pdf"), conn)
    assert result.status == "ingested" and result.figures == 1
    row = conn.execute("SELECT caption, description FROM figures").fetchone()
    assert row["description"] is None and row["caption"].startswith("Figure 3")
    assert _count(conn, "figure_vectors") == 1  # caption alone is still searchable
    import json

    flags = json.loads(conn.execute("SELECT gap_flags FROM documents").fetchone()["gap_flags"])
    assert any("figure description failed" in f for f in flags)


def test_reingest_inherits_category_and_source_uri_from_raw_store(tmp_path, conn, fake_passes, monkeypatch):
    """A raw-store re-ingest must not wipe the facets: category and source_uri inherit
    from the replaced doc (observed live: a 'paper' became 'uncategorized')."""
    from locus.config import load as cfg_load

    raw_store = tmp_path / "raw"
    raw_store.mkdir()
    monkeypatch.setattr(cfg_load().paths, "raw_store", raw_store)

    pdf = _make_pdf(tmp_path / "doc.pdf")
    first = ingest_file(pdf, conn)
    conn.execute(
        "UPDATE documents SET category='paper', source_uri='vault/incoming/papers/orig.pdf' "
        "WHERE id=?", (first.doc_id,),
    )
    conn.commit()

    # Re-ingest the identical bytes from inside the raw store (the ONE-re-ingest path).
    import shutil as _sh

    raw_copy = raw_store / "deadbeef.pdf"
    _sh.copy2(pdf, raw_copy)
    result = ingest_file(raw_copy, conn, reingest=True)
    assert result.status == "ingested"
    row = conn.execute("SELECT category, source_uri FROM documents").fetchone()
    assert row["category"] == "paper"  # inherited, not 'uncategorized'
    assert row["source_uri"] == "vault/incoming/papers/orig.pdf"  # original location kept

    # A re-drop from a NON-raw-store path is a genuine new source: uri updates, category
    # still inherits when the new path derives none.
    result2 = ingest_file(pdf, conn, reingest=True)
    row2 = conn.execute("SELECT category, source_uri FROM documents").fetchone()
    assert result2.status == "ingested"
    assert row2["category"] == "paper"
    assert row2["source_uri"] == str(pdf)
