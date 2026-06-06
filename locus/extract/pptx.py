"""Extract structured content from a .pptx file with python-pptx.

Produces an `ExtractedDoc` (see extract/base.py). Each slide yields one (title, text,
slide_no) span: title from the title placeholder, body from every text-bearing shape
(except slide chrome, filtered two ways: slide-number / footer / date *placeholders*,
and plain text boxes whose entire text is the slide's own number — how Google-Slides
exports render page numbers), tables
flattened to pipe-joined rows in-flow, and speaker notes appended under a
"Notes:" marker. Spans are then size-banded like every other format — but unlike the
unpaginated formats, sections carry REAL slide-number ranges in page_start/page_end
(a slide IS a page), so citations and section_map point at slides.

LIMITATIONS (documented, not bugs):
- Shapes are taken in XML order, not geometric reading order. Placeholders commonly
  inherit top/left from the slide layout (None on the shape), so a geometric sort would
  itself be a guess. Acceptable: retrieval grounds on chunks regardless of intra-section
  ordering (§15.0) — same tradeoff docx.py takes for tables.
- Hidden slides are included (python-pptx has no first-class hidden flag); authored
  content stays retrievable.
- Slide images / charts / SmartArt contribute no text. With `figures=True` (the pipeline's
  setting), visual-bearing slides are instead captured WHOLE: the deck is converted to PDF
  via LibreOffice (`soffice --headless`) and each qualifying slide page is rendered to PNG
  as an ExtractedFigure(kind="slide") — the step-11 fix for chart-blind decks. Requires
  `soffice` on PATH; absent (or conversion failure / page-count mismatch from hidden
  slides) the deck degrades gracefully: text extraction proceeds, no figures, the reason
  recorded in `figure_fallbacks` (-> gap_flags). Never quarantines (§6).

Section text includes its own heading line (the title placeholder falls out of the
shape walk), mirroring pdf.py / docx.py / textdoc.py.
"""

from __future__ import annotations

import logging
import shutil
import subprocess
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from locus.extract.base import (
    MIN_SECTION_CHARS,
    ExtractedDoc,
    ExtractedFigure,
    ExtractedSection,
    has_math,
    window_by_chars,
)

log = logging.getLogger(__name__)

_NOTES_MARKER = "Notes:"

# Page furniture, not content: slide-number / footer / date placeholders pollute every
# slide's text with stray numbers and boilerplate (observed on the first real deck —
# chunks led with the literal slide number). PowerPoint injects these into slide XML
# when "Insert > Header & Footer" options are on.
_CHROME_PLACEHOLDERS = {
    PP_PLACEHOLDER.SLIDE_NUMBER,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.DATE,
}

# (title, text, slide_no) — slide_no is 1-based
_SlideSpan = tuple[str | None, str, int]


def extract_pptx(path: str | Path, *, figures: bool = False) -> ExtractedDoc:
    """Extract a structured ExtractedDoc from the .pptx file at `path`.

    `figures=True` (the pipeline's setting) renders visual-bearing slides to PNG via
    LibreOffice and attaches them as `doc.figures` (see module docstring). Default False so
    ad-hoc extraction never shells out to soffice.
    """
    path = Path(path)
    try:
        prs = Presentation(str(path))
    except Exception as exc:  # python-pptx raises several types on a damaged/non-pptx file
        raise ValueError(f"not a readable pptx: {exc}") from exc

    slide_count = len(prs.slides)
    spans = _slide_spans(prs)
    if not any(text.strip() for _, text, _ in spans):
        raise ValueError("no extractable text")

    sections = _build_slide_sections(spans)
    if any(title is not None for title, _, _ in spans):
        strategy = "slides"
    elif len(sections) > 1:
        strategy = "windowed"
    else:
        strategy = "single"

    figure_list: list[ExtractedFigure] = []
    figure_fallbacks: list[str] = []
    if figures:
        from locus.config import load  # lazy, mirrors pdf.py

        fig_cfg = load().figures
        if fig_cfg.enabled and fig_cfg.render_slides:
            figure_list, figure_fallbacks = _slide_figures(path, prs)
            from locus.extract.figures_detect import assign_sections

            assign_sections(figure_list, sections)

    return ExtractedDoc(
        title=_resolve_title(prs, spans, path),
        page_count=slide_count,  # a slide IS a page; sections carry real slide ranges
        section_strategy=strategy,
        sections=sections,
        source_path=str(path),
        source_date=_resolve_source_date(prs),
        figures=figure_list,
        figure_fallbacks=figure_fallbacks,
    )


def _slide_spans(prs) -> list[_SlideSpan]:
    """One (title, text, slide_no) span per slide; truly empty slides are skipped."""
    spans: list[_SlideSpan] = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        body = _slide_body(slide, slide_no)
        notes = _slide_notes(slide)
        parts = [p for p in (body, f"{_NOTES_MARKER}\n{notes}" if notes else "") if p]
        text = "\n\n".join(parts)
        if not text.strip():
            continue
        spans.append((title, text, slide_no))
    return spans


def _slide_title(slide) -> str | None:
    """Text of the title placeholder, or None when absent/empty."""
    shape = slide.shapes.title
    if shape is None:
        return None
    return shape.text_frame.text.strip() or None


def _slide_body(slide, slide_no: int) -> str:
    """All text-bearing shapes in XML order; the title placeholder leads naturally.
    Slide chrome is skipped: placeholder furniture, and literal slide-number text boxes
    (Google-Slides-style exports render the page number as a plain TEXT_BOX whose whole
    text is the slide number — observed on the first real deck)."""
    blocks: list[str] = []
    for shape in slide.shapes:
        if _is_chrome(shape):
            continue
        text = _shape_text(shape)
        if text.strip() == str(slide_no):
            continue  # a shape whose entire text is this slide's own number is chrome
        if text.strip():
            blocks.append(text.strip())
    return "\n\n".join(blocks)


def _is_chrome(shape) -> bool:
    """Slide-number / footer / date placeholder — boilerplate, never slide content."""
    return shape.is_placeholder and shape.placeholder_format.type in _CHROME_PLACEHOLDERS


def _shape_text(shape) -> str:
    """Text frame content, or pipe-joined table rows for graphic-frame tables."""
    if shape.has_text_frame:
        return shape.text_frame.text
    if getattr(shape, "has_table", False):  # has_table exists on GraphicFrame only
        return _table_as_text(shape.table)
    return ""  # pictures, charts, media — no text layer (figures are block-11 work)


def _table_as_text(table) -> str:
    """A table flattened to pipe-joined rows (docx.py precedent)."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def _slide_notes(slide) -> str:
    """Speaker-notes body text, or "". Gated on has_notes_slide: accessing .notes_slide
    on a slide without one would CREATE it in the in-memory package. notes_text_frame is
    the notes-body placeholder specifically — never the slide-number placeholder."""
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def _build_slide_sections(spans: list[_SlideSpan]) -> list[ExtractedSection]:
    """Size-banded sections carrying slide-number ranges.

    base.build_sections with slide ranges threaded through: merge consecutive under-MIN
    spans (range = first..last absorbed slide; a merged section keeps the first non-empty
    heading — span text is verbatim and includes its own heading line, so absorbed
    headings stay visible), window over-MAX spans (window parts inherit the merged
    group's slide range — a single slide over MAX is rare enough that the inherited
    range stays honest), then assign positions and the has_math flag.
    """
    spans = [(t, x, n) for t, x, n in spans if x.strip()]

    # (title, text, slide_start, slide_end)
    merged: list[tuple[str | None, str, int, int]] = []
    cur_title: str | None = None
    cur_text = ""
    cur_start = cur_end = 0
    for title, text, slide_no in spans:
        if not cur_text:
            cur_title, cur_text, cur_start = title, text, slide_no
        else:
            cur_text += "\n\n" + text
            if cur_title is None:
                cur_title = title
        cur_end = slide_no
        if len(cur_text.strip()) >= MIN_SECTION_CHARS:
            merged.append((cur_title, cur_text, cur_start, cur_end))
            cur_title, cur_text = None, ""
    if cur_text:  # trailing under-sized bucket: fold into the previous section if any
        if merged and len(cur_text.strip()) < MIN_SECTION_CHARS:
            pt, px, ps, _ = merged[-1]
            merged[-1] = (pt, px + "\n\n" + cur_text, ps, cur_end)
        else:
            merged.append((cur_title, cur_text, cur_start, cur_end))

    windowed: list[tuple[str | None, str, int, int]] = []
    for title, text, start, end in merged:
        for w_title, w_text in window_by_chars(title, text):
            windowed.append((w_title, w_text, start, end))

    return [
        ExtractedSection(
            position=i,
            title=title,
            text=text.strip(),
            page_start=start,
            page_end=end,
            has_math=has_math(text),
        )
        for i, (title, text, start, end) in enumerate(windowed)
    ]


# --- slide figures (step 11): full-slide render of visual-bearing slides ------------------

# Shape types that ARE visual content on their own, at any size.
_ALWAYS_VISUAL = {
    MSO_SHAPE_TYPE.CHART,
    MSO_SHAPE_TYPE.DIAGRAM,  # SmartArt
    MSO_SHAPE_TYPE.MEDIA,
}
# Pictures are visual only when sizeable: every slide of a templated deck carries a small
# logo PICTURE (measured 0.26% of slide area on the first real deck), while content
# pictures/screenshots run far larger (56%). The [figures] min_area_frac threshold (3%)
# separates them cleanly. AUTO_SHAPE/FREEFORM counts are deliberately NOT a signal:
# consulting-style decks build text layouts from large autoshape cards (measured 13-29%
# of slide area each), indistinguishable from drawn diagrams by shape census — a
# PowerPoint-drawn box-and-arrow diagram is therefore a documented miss (its box text is
# still extracted; only the arrows are lost), the err-strict side of the tradeoff.
_PICTURE_TYPES = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}
_SOFFICE_TIMEOUT = 180.0  # seconds; whole-deck PDF conversion
_RENDER_ZOOM = 2.0  # matches figures_detect: enough detail for the VLM


def _slide_has_visual(slide, slide_area: int, min_area_frac: float) -> bool:
    """True when a slide carries visual content the text layer cannot represent."""

    def walk(shapes) -> bool:
        for shape in shapes:
            try:
                stype = shape.shape_type
            except Exception:  # python-pptx raises on some exotic shape XML
                continue
            if stype in _ALWAYS_VISUAL or getattr(shape, "has_chart", False):
                return True
            if stype in _PICTURE_TYPES and slide_area > 0:
                area = (shape.width or 0) * (shape.height or 0)
                if area / slide_area >= min_area_frac:
                    return True
            if stype == MSO_SHAPE_TYPE.GROUP and walk(shape.shapes):
                return True
        return False

    return walk(slide.shapes)


def _slide_figures(path: Path, prs) -> tuple[list[ExtractedFigure], list[str]]:
    """Render each visual-bearing slide to a PNG ExtractedFigure(kind='slide').

    Whole-deck conversion via `soffice --headless --convert-to pdf`, then per-page render
    with pymupdf. The slide->page mapping is trusted only when the rendered page count
    equals the slide count (hidden slides can be dropped from the export, silently
    shifting every page); on mismatch — or soffice missing/failing — the deck degrades
    to no figures with the reason returned for the audit trail. Never raises.
    """
    from locus.config import load

    min_area_frac = load().figures.min_area_frac
    slide_area = prs.slide_width * prs.slide_height
    visual = [
        no
        for no, slide in enumerate(prs.slides, start=1)
        if _slide_has_visual(slide, slide_area, min_area_frac)
    ]
    if not visual:
        return [], []
    if shutil.which("soffice") is None:
        log.warning("slide figures unavailable for %s: soffice (LibreOffice) not on PATH", path)
        return [], ["slide figures unavailable: soffice (LibreOffice) not on PATH"]

    import pymupdf  # local import keeps the soffice-less path import-free

    try:
        with tempfile.TemporaryDirectory(prefix="locus-slides-") as td:
            proc = subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", td, str(path)],
                capture_output=True,
                timeout=_SOFFICE_TIMEOUT,
            )
            rendered = Path(td) / f"{path.stem}.pdf"
            if proc.returncode != 0 or not rendered.exists():
                reason = (proc.stderr or proc.stdout or b"").decode(errors="replace").strip()
                log.warning("soffice conversion failed for %s: %s", path, reason[:200])
                return [], ["slide figures unavailable: soffice conversion failed"]
            doc = pymupdf.open(rendered)
            try:
                if doc.page_count != len(prs.slides):
                    log.warning(
                        "slide figures skipped for %s: rendered %d pages != %d slides "
                        "(hidden slides?)", path, doc.page_count, len(prs.slides),
                    )
                    return [], [
                        f"slide figures unavailable: rendered {doc.page_count} pages "
                        f"!= {len(prs.slides)} slides"
                    ]
                figures: list[ExtractedFigure] = []
                for slide_no in visual:
                    png = doc[slide_no - 1].get_pixmap(
                        matrix=pymupdf.Matrix(_RENDER_ZOOM, _RENDER_ZOOM)
                    ).tobytes("png")
                    figures.append(
                        ExtractedFigure(
                            page=slide_no,
                            image_bytes=png,
                            caption=_slide_caption(prs.slides[slide_no - 1], slide_no),
                            kind="slide",
                        )
                    )
                return figures, []
            finally:
                doc.close()
    except subprocess.TimeoutExpired:
        log.warning("soffice conversion timed out for %s", path)
        return [], ["slide figures unavailable: soffice conversion timed out"]
    except Exception as exc:  # degradation, never quarantine (§6)
        log.warning("slide figure render failed for %s: %s", path, exc)
        return [], [f"slide figures unavailable: {exc}"]


_CAPTION_MAX_CHARS = 160  # a caption is the slide's heading line, not its body


def _slide_caption(slide, slide_no: int) -> str | None:
    """The slide's visible heading, for use as the rendered figure's caption.

    Title placeholder when present; otherwise the first non-empty body line — decks
    exported without placeholder structure (the first real deck: Google-Slides export,
    zero title placeholders) put the visible heading in a plain text box, which leads
    the body text. Deterministic; None when the slide has no text at all.
    """
    title = _slide_title(slide)
    if title:
        return title
    for line in _slide_body(slide, slide_no).splitlines():
        line = line.strip()
        if line:
            return line[:_CAPTION_MAX_CHARS]
    return None


def _resolve_title(prs, spans: list[_SlideSpan], path: Path) -> str | None:
    """Core-properties title, else the first slide title, else the filename stem."""
    meta = (prs.core_properties.title or "").strip()
    if len(meta) > 3:
        return meta
    for title, _, _ in spans:
        if title:
            return title
    return path.stem


def _resolve_source_date(prs) -> str | None:
    """Core-properties created date if present, else modified date, else None."""
    props = prs.core_properties
    for stamp in (props.created, props.modified):
        if stamp is not None:
            return stamp.date().isoformat()
    return None
